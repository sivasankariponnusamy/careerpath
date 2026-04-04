from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import pandas as pd
import numpy as np
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.preprocessing import StandardScaler, LabelEncoder
from sklearn.model_selection import train_test_split
import os
import re
import json
import requests
from werkzeug.utils import secure_filename
from datetime import datetime
import shutil

# Import document readers with fallbacks
try:
    from PyPDF2 import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

app = Flask(__name__)
CORS(app)

# Database Configuration
from models import db, Resume, SkillGapAnalysis

BASEDIR = os.path.abspath(os.path.dirname(__file__))
db_path = os.environ.get('DB_PATH', os.path.join(BASEDIR, 'career_path.db'))
upload_folder = os.environ.get('UPLOAD_FOLDER', os.path.join(BASEDIR, 'uploads'))
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + db_path
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = upload_folder
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Initialize database
db.init_app(app)

# Create uploads directory if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Create database tables
with app.app_context():
    db.create_all()
    print("✓ Database initialized successfully!")

# LLaMA Configuration (using Ollama)
OLLAMA_API_URL = "http://localhost:11434/api/generate"
LLAMA_MODEL = "llama2"  # You can change this to llama3 or other models

# Load datasets
DATASET_PATH = os.path.join(os.path.dirname(__file__), '..', 'dataset')
try:
    df_text = pd.read_csv(os.path.join(DATASET_PATH, 'job_dataset_text.csv'))
    df_encoded = pd.read_csv(os.path.join(DATASET_PATH, 'job_dataset_encoded.csv'))
    all_skills = [col for col in df_encoded.columns if col != 'role']
    print(f"✓ Dataset loaded with {len(all_skills)} skills")
except Exception as e:
    print(f"⚠ Dataset not found, using built-in skill list: {e}")
    df_text = pd.DataFrame()
    df_encoded = pd.DataFrame()
    # Comprehensive built-in skill list for serverless deployment
    all_skills = [
        'Python', 'JavaScript', 'TypeScript', 'Java', 'C++', 'C#', 'Go', 'Rust', 'PHP', 'Ruby',
        'Swift', 'Kotlin', 'Scala', 'R', 'MATLAB', 'Perl', 'Bash', 'Shell',
        'React', 'Angular', 'Vue', 'HTML', 'CSS', 'Node.js', 'Express', 'Flask', 'Django',
        'FastAPI', 'Spring', 'Spring Boot', 'Laravel', 'Rails', 'ASP.NET', 'Next.js', 'Nuxt.js',
        'Tailwind', 'Bootstrap', 'jQuery', 'Redux', 'GraphQL', 'REST API', 'gRPC',
        'MongoDB', 'PostgreSQL', 'MySQL', 'SQLite', 'SQL', 'NoSQL', 'Redis', 'Cassandra',
        'Oracle', 'DynamoDB', 'Elasticsearch', 'Firebase', 'Supabase',
        'Docker', 'Kubernetes', 'AWS', 'Azure', 'GCP', 'CI/CD', 'Jenkins', 'GitHub Actions',
        'Git', 'Linux', 'Nginx', 'Apache', 'Terraform', 'Ansible', 'Prometheus', 'Grafana',
        'EC2', 'S3', 'Lambda', 'CloudFormation', 'EKS', 'ECS', 'RDS',
        'Machine Learning', 'Deep Learning', 'TensorFlow', 'PyTorch', 'Scikit-learn',
        'NLP', 'Computer Vision', 'OpenCV', 'LangChain', 'RAG', 'LLM',
        'Pandas', 'NumPy', 'Matplotlib', 'Seaborn', 'Jupyter', 'Data Analysis',
        'PowerBI', 'Tableau', 'Excel', 'Data Engineering', 'Apache Spark', 'Kafka',
        'Android', 'iOS', 'React Native', 'Flutter', 'Dart', 'Xcode',
        'Figma', 'Adobe XD', 'Sketch', 'UI/UX', 'Wireframing', 'Prototyping',
        'Agile', 'Scrum', 'Jira', 'Project Management', 'Leadership',
        'Problem Solving', 'Teamwork', 'Communication', 'Critical Thinking',
        'Cybersecurity', 'Network Security', 'Penetration Testing', 'OWASP',
        'Blockchain', 'Solidity', 'Web3', 'Smart Contracts',
        'OpenAI', 'HuggingFace', 'BERT', 'GPT', 'Stable Diffusion',
        'RabbitMQ', 'Celery', 'WebSocket', 'OAuth', 'JWT', 'Microservices',
        'Unit Testing', 'Jest', 'Pytest', 'Selenium', 'Cypress',
    ]

# If dataset was not loaded, all_skills is already set above; otherwise use dataset columns
if not df_encoded.empty:
    all_skills = [col for col in df_encoded.columns if col != 'role']

class CareerRecommendationSystem:
    """
    AI-Driven Career Recommendation System using Random Forest Classifier
    """
    def __init__(self):
        self.rf_model = None
        self.label_encoder = LabelEncoder()
        self.scaler = StandardScaler()
        self.role_profiles = self._create_role_profiles()
        self.train_model()
    
    def _create_role_profiles(self):
        """Create average skill profiles for each role"""
        if df_encoded.empty:
            return {}
        profiles = {}
        for role in df_encoded['role'].unique():
            role_data = df_encoded[df_encoded['role'] == role]
            avg_profile = role_data[all_skills].mean().values
            profiles[role] = avg_profile
        return profiles
    
    def train_model(self):
        """Train Random Forest Classifier on the dataset"""
        if df_encoded.empty:
            print("⚠ No dataset available, skipping model training")
            return
        print("Training Random Forest Classifier...")
        
        # Prepare features and labels
        X = df_encoded[all_skills].values
        y = df_encoded['role'].values
        
        # Encode labels
        y_encoded = self.label_encoder.fit_transform(y)
        
        # Split data for training
        X_train, X_test, y_train, y_test = train_test_split(
            X, y_encoded, test_size=0.2, random_state=42, stratify=y_encoded
        )
        
        # Train Random Forest Classifier
        self.rf_model = RandomForestClassifier(
            n_estimators=100,
            max_depth=20,
            min_samples_split=5,
            min_samples_leaf=2,
            random_state=42,
            n_jobs=-1
        )
        
        self.rf_model.fit(X_train, y_train)
        
        # Calculate accuracy
        train_accuracy = self.rf_model.score(X_train, y_train)
        test_accuracy = self.rf_model.score(X_test, y_test)
        
        print(f"✓ Model trained successfully!")
        print(f"  Training Accuracy: {train_accuracy * 100:.2f}%")
        print(f"  Testing Accuracy: {test_accuracy * 100:.2f}%")
    
    def predict_career(self, user_skills):
        """Predict career role using Random Forest Classifier"""
        if self.rf_model is None or not all_skills:
            # Fallback: rule-based role detection when no ML model
            return self._rule_based_predict(user_skills)

        # Create user skill vector
        user_vector = np.zeros(len(all_skills))
        for skill in user_skills:
            if skill in all_skills:
                idx = all_skills.index(skill)
                user_vector[idx] = 1
        
        # Reshape for prediction
        user_vector = user_vector.reshape(1, -1)
        
        # Get prediction probabilities
        probabilities = self.rf_model.predict_proba(user_vector)[0]
        
        # Get top predictions
        top_indices = np.argsort(probabilities)[::-1][:5]
        predictions = []
        
        for idx in top_indices:
            role = self.label_encoder.inverse_transform([idx])[0]
            confidence = probabilities[idx] * 100
            predictions.append({
                'role': role,
                'confidence': round(confidence, 2),
                'match_percentage': round(confidence, 2)
            })
        
        return predictions

    def _rule_based_predict(self, user_skills):
        """Rule-based role prediction when ML model is not available"""
        skills_lower = [s.lower() for s in user_skills]
        role_scores = {
            'Frontend Developer': sum(1 for s in ['react', 'angular', 'vue', 'html', 'css', 'javascript', 'typescript', 'tailwind', 'next.js'] if s in skills_lower),
            'Backend Developer': sum(1 for s in ['python', 'java', 'node.js', 'flask', 'django', 'fastapi', 'spring', 'postgresql', 'mysql', 'rest api'] if s in skills_lower),
            'Full Stack Developer': sum(1 for s in ['react', 'node.js', 'python', 'javascript', 'mongodb', 'postgresql', 'html', 'css'] if s in skills_lower),
            'Data Scientist': sum(1 for s in ['python', 'machine learning', 'pandas', 'numpy', 'scikit-learn', 'tensorflow', 'pytorch', 'data analysis', 'r'] if s in skills_lower),
            'DevOps Engineer': sum(1 for s in ['docker', 'kubernetes', 'aws', 'azure', 'ci/cd', 'linux', 'terraform', 'jenkins', 'git'] if s in skills_lower),
            'Mobile Developer': sum(1 for s in ['android', 'ios', 'react native', 'flutter', 'swift', 'kotlin', 'dart'] if s in skills_lower),
            'Machine Learning Engineer': sum(1 for s in ['machine learning', 'deep learning', 'tensorflow', 'pytorch', 'scikit-learn', 'nlp', 'computer vision', 'python'] if s in skills_lower),
            'Data Engineer': sum(1 for s in ['python', 'apache spark', 'kafka', 'sql', 'aws', 'data engineering', 'pandas'] if s in skills_lower),
            'Cloud Engineer': sum(1 for s in ['aws', 'azure', 'gcp', 'docker', 'kubernetes', 'terraform', 'ec2', 's3'] if s in skills_lower),
            'UI/UX Designer': sum(1 for s in ['figma', 'adobe xd', 'sketch', 'ui/ux', 'wireframing', 'prototyping', 'css'] if s in skills_lower),
        }
        sorted_roles = sorted(role_scores.items(), key=lambda x: x[1], reverse=True)
        predictions = []
        for role, score in sorted_roles[:5]:
            if score > 0:
                confidence = min(score * 15, 95)
                predictions.append({'role': role, 'confidence': confidence, 'match_percentage': confidence})
        return predictions if predictions else [{'role': 'Software Developer', 'confidence': 50, 'match_percentage': 50}]
    
    def analyze_skill_gap(self, user_skills, target_role):
        """Analyze skill gap between user skills and target role"""
        # Create user skill vector
        user_vector = np.zeros(len(all_skills))
        for skill in user_skills:
            if skill in all_skills:
                idx = all_skills.index(skill)
                user_vector[idx] = 1
        
        # Get target role profile
        if target_role not in self.role_profiles:
            return None
        
        target_vector = self.role_profiles[target_role]
        
        # Calculate similarity
        similarity = cosine_similarity([user_vector], [target_vector])[0][0]
        match_percentage = similarity * 100
        
        # Find missing skills (required but not possessed)
        missing_skills = []
        for i, skill in enumerate(all_skills):
            if user_vector[i] == 0 and target_vector[i] > 0.3:  # Threshold for important skills
                missing_skills.append(skill)
        
        # Find matching skills
        matching_skills = []
        for i, skill in enumerate(all_skills):
            if user_vector[i] == 1 and target_vector[i] > 0.2:
                matching_skills.append(skill)
        
        # Find recommended skills (high importance in role profile)
        recommended_skills = []
        for i, skill in enumerate(all_skills):
            if target_vector[i] > 0.5:  # High importance threshold
                recommended_skills.append({
                    'skill': skill,
                    'importance': round(target_vector[i] * 100, 2),
                    'has_skill': bool(user_vector[i])
                })
        
        recommended_skills.sort(key=lambda x: x['importance'], reverse=True)
        
        return {
            'match_percentage': round(match_percentage, 2),
            'missing_skills': missing_skills,
            'matching_skills': matching_skills,
            'recommended_skills': recommended_skills[:10],  # Top 10 most important skills
            'total_skills_needed': int(np.sum(target_vector > 0.3)),
            'skills_you_have': len(matching_skills)
        }
    
    def recommend_roles(self, user_skills, top_n=5):
        """Recommend top matching roles using ML model predictions"""
        return self.predict_career(user_skills)[:top_n]

class LLaMAIntegration:
    """
    LLaMA AI Integration for Intelligent Skill-Gap Explanation
    and Learning Roadmap Generation
    """
    
    @staticmethod
    def check_ollama_status():
        """Check if Ollama is running"""
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            return response.status_code == 200
        except:
            return False
    
    @staticmethod
    def generate_response(prompt, max_tokens=1000):
        """Generate response using LLaMA via Ollama"""
        try:
            payload = {
                "model": LLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.7,
                    "num_predict": max_tokens
                }
            }
            
            response = requests.post(
                OLLAMA_API_URL,
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get('response', '').strip()
            else:
                return None
                
        except Exception as e:
            print(f"LLaMA API Error: {str(e)}")
            return None
    
    @staticmethod
    def explain_skill_gap(user_skills, target_role, missing_skills, matching_skills):
        """Generate intelligent skill-gap explanation using LLaMA"""
        
        prompt = f"""You are an expert career counselor for Computer Science students.

Student's Current Skills: {', '.join(user_skills)}
Target Career Role: {target_role}
Skills Already Possessed: {', '.join(matching_skills)}
Missing Skills: {', '.join(missing_skills)}

Task: Provide a brief, personalized skill-gap analysis and encouragement (maximum 150 words). 
Focus on:
1. Acknowledge their existing strengths
2. Explain why the missing skills are important for {target_role}
3. Provide motivation and next steps

Keep the tone professional yet encouraging."""

        response = LLaMAIntegration.generate_response(prompt, max_tokens=250)
        
        if response:
            return response
        else:
            # Fallback response if LLaMA is unavailable
            return f"""Based on your skill profile, you have a strong foundation with skills like {', '.join(matching_skills[:3])}. To excel as a {target_role}, you'll need to develop {len(missing_skills)} additional skills including {', '.join(missing_skills[:3])}. These skills are essential for this role and will significantly enhance your career prospects. Focus on systematic learning and hands-on practice."""
    
    @staticmethod
    def generate_learning_roadmap(user_skills, target_role, missing_skills):
        """Generate 6-month learning roadmap using LLaMA"""
        
        prompt = f"""You are an expert career mentor for Computer Science students.

Student's Skills: {', '.join(user_skills)}
Target Role: {target_role}
Skills to Learn: {', '.join(missing_skills[:8])}

Task: Create a structured 6-month learning roadmap with 3 phases:
- Phase 1 (Month 1-2): Foundation skills
- Phase 2 (Month 3-4): Intermediate skills  
- Phase 3 (Month 5-6): Advanced skills and projects

For each phase, specify:
1. Which skills to focus on
2. Recommended learning resources
3. Practical project ideas

Keep it actionable and realistic. Maximum 300 words."""

        response = LLaMAIntegration.generate_response(prompt, max_tokens=500)
        
        if response:
            return response
        else:
            # Fallback structured roadmap
            priority_skills = missing_skills[:6] if len(missing_skills) >= 6 else missing_skills
            
            roadmap = {
                'phase1': {
                    'duration': 'Month 1-2: Foundation',
                    'skills': priority_skills[:2] if len(priority_skills) >= 2 else priority_skills,
                    'focus': 'Build strong fundamentals through online courses and tutorials',
                    'resources': ['Coursera', 'Udemy', 'freeCodeCamp']
                },
                'phase2': {
                    'duration': 'Month 3-4: Intermediate',
                    'skills': priority_skills[2:4] if len(priority_skills) >= 4 else [],
                    'focus': 'Apply knowledge through practical projects',
                    'resources': ['YouTube tutorials', 'Documentation', 'GitHub projects']
                },
                'phase3': {
                    'duration': 'Month 5-6: Advanced',
                    'skills': priority_skills[4:] if len(priority_skills) > 4 else [],
                    'focus': 'Build portfolio projects and contribute to open source',
                    'resources': ['GitHub', 'Kaggle', 'HackerRank', 'Personal projects']
                }
            }
            
            return roadmap

# Initialize systems
career_system = CareerRecommendationSystem()
llama_ai = LLaMAIntegration()

@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    ollama_status = llama_ai.check_ollama_status()
    return jsonify({
        'status': 'healthy',
        'message': 'Flask backend is running',
        'ml_model': 'Random Forest Classifier',
        'ml_model_trained': career_system.rf_model is not None,
        'llama_available': ollama_status,
        'llama_model': LLAMA_MODEL if ollama_status else 'Not available'
    })

@app.route('/api/skills', methods=['GET'])
def get_all_skills():
    """Get all available skills"""
    return jsonify({'skills': all_skills, 'count': len(all_skills)})

@app.route('/api/extract-skills', methods=['POST'])
def extract_skills_from_resume():
    """Extract skills from resume (PDF, Word, PowerPoint, or text files)"""
    
    if 'resume' not in request.files:
        return jsonify({'error': 'No resume file provided'}), 400
    
    file = request.files['resume']
    original_filename = file.filename
    filename = file.filename.lower()
    
    if not original_filename:
        return jsonify({'error': 'No file selected'}), 400
    
    # Save the file first
    secure_fname = secure_filename(original_filename)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    unique_filename = f"{timestamp}_{secure_fname}"
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
    
    # Get file size
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    
    # Save file to disk
    file.save(file_path)
    
    # Read file content for processing
    file_type = filename.split('.')[-1] if '.' in filename else 'unknown'
    content = ''
    
    try:
        if filename.endswith('.pdf'):
            if not PDF_SUPPORT:
                return jsonify({'error': 'PDF support not available. Please install PyPDF2.'}), 500
            
            # Read PDF content
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                content += page.extract_text() + '\n'
                
        elif filename.endswith('.docx'):
            if not DOCX_SUPPORT:
                return jsonify({'error': 'DOCX support not available. Please install python-docx.'}), 500
            
            # Read DOCX content
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                content += paragraph.text + '\n'
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        content += cell.text + ' '
                content += '\n'
                
        elif filename.endswith('.pptx'):
            if not PPTX_SUPPORT:
                return jsonify({'error': 'PPTX support not available. Please install python-pptx.'}), 500
            
            # Read PPTX content
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'
                        
        elif filename.endswith(('.txt', '.text')):
            # Handle as plain text file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
        else:
            # Try to read as text for any other format
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            except:
                # Clean up file and return error
                if os.path.exists(file_path):
                    os.remove(file_path)
                return jsonify({'error': f'Unsupported file format: {filename}. Please upload PDF, DOCX, PPTX, or TXT files.'}), 400
                
    except Exception as e:
        # Clean up file on error
        if os.path.exists(file_path):
            os.remove(file_path)
        return jsonify({'error': f'Could not read file: {str(e)}'}), 400
    
    if not content.strip():
        # Clean up file and return error
        if os.path.exists(file_path):
            os.remove(file_path)
        return jsonify({'error': 'No text content found in the file'}), 400
    
    # Extract skills by matching against our skill list
    # Convert content to lowercase for case-insensitive matching
    content_lower = content.lower()
    
    # Find skills that appear in the resume
    found_skills = []
    for skill in all_skills:
        # Create a pattern that matches the skill as a whole word
        pattern = r'\b' + re.escape(skill.lower()) + r'\b'
        if re.search(pattern, content_lower):
            found_skills.append(skill)
    
    # Also check for common programming languages and technologies
    # that might be formatted differently in resumes
    skill_aliases = {
        'Python': ['python', 'py'],
        'JavaScript': ['javascript', 'js', 'node.js', 'nodejs'],
        'TypeScript': ['typescript', 'ts'],
        'Java': ['java'],
        'C++': ['c++', 'cpp'],
        'C#': ['c#', 'csharp'],
        'Go': ['go', 'golang'],
        'Rust': ['rust'],
        'React': ['react', 'reactjs', 'react.js'],
        'Angular': ['angular', 'angularjs'],
        'Vue': ['vue', 'vuejs', 'vue.js'],
        'Flask': ['flask'],
        'Django': ['django'],
        'MongoDB': ['mongodb', 'mongo'],
        'PostgreSQL': ['postgresql', 'postgres'],
        'MySQL': ['mysql'],
        'Docker': ['docker'],
        'Kubernetes': ['kubernetes', 'k8s'],
        'AWS': ['aws', 'amazon web services'],
        'Azure': ['azure', 'microsoft azure'],
        'Machine Learning': ['machine learning', 'ml'],
        'Deep Learning': ['deep learning', 'dl'],
        'TensorFlow': ['tensorflow'],
        'PyTorch': ['pytorch'],
        'Scikit-learn': ['scikit-learn', 'sklearn'],
        'HTML': ['html', 'html5'],
        'CSS': ['css', 'css3'],
        'REST API': ['rest api', 'restful', 'rest'],
        'Git': ['git', 'github', 'gitlab'],
        'Linux': ['linux', 'unix'],
        'SQL': ['sql'],
        'NoSQL': ['nosql'],
        'Data Analysis': ['data analysis', 'data analytics'],
        'Problem Solving': ['problem solving', 'problem-solving'],
        'Project Management': ['project management'],
        'Teamwork': ['teamwork', 'team collaboration'],
        'OpenCV': ['opencv'],
        'NLP': ['nlp', 'natural language processing'],
        'Computer Vision': ['computer vision', 'cv'],
        'LangChain': ['langchain'],
        'RAG': ['rag', 'retrieval augmented generation'],
        'Elasticsearch': ['elasticsearch', 'elastic search'],
        'Tailwind': ['tailwind', 'tailwindcss'],
        'EC2': ['ec2'],
        'S3': ['s3'],
        'PowerBI': ['powerbi', 'power bi'],
    }
    
    # Match skills using aliases
    matched_skills = set()
    for canonical_skill, aliases in skill_aliases.items():
        for alias in aliases:
            pattern = r'\b' + re.escape(alias.lower()) + r'\b'
            if re.search(pattern, content_lower):
                # Check if this skill exists in our dataset
                if canonical_skill in all_skills:
                    matched_skills.add(canonical_skill)
                # Also add the exact match if it exists
                for skill in all_skills:
                    if skill.lower() == canonical_skill.lower():
                        matched_skills.add(skill)
                break
    
    # Combine found skills and remove duplicates
    all_found_skills = list(set(found_skills + list(matched_skills)))
    
    # Categorize skills
    skill_categories = {
        'Programming Languages': ['Python', 'JavaScript', 'TypeScript', 'Java', 'C++', 'C#', 'Go', 'Rust', 'PHP', 'Ruby', 'Swift', 'Kotlin'],
        'Web Technologies': ['React', 'Angular', 'Vue', 'HTML', 'CSS', 'Flask', 'Django', 'Node.js', 'Express', 'FastAPI', 'Tailwind'],
        'Databases': ['MongoDB', 'PostgreSQL', 'MySQL', 'SQL', 'NoSQL', 'Redis', 'Cassandra', 'Oracle'],
        'DevOps & Cloud': ['Docker', 'Kubernetes', 'AWS', 'Azure', 'GCP', 'CI/CD', 'Jenkins', 'Git', 'Linux', 'EC2', 'S3'],
        'AI & Machine Learning': ['Machine Learning', 'Deep Learning', 'TensorFlow', 'PyTorch', 'Scikit-learn', 'NLP', 'Computer Vision', 'OpenCV', 'LangChain', 'RAG'],
        'APIs & Tools': ['REST API', 'GraphQL', 'Elasticsearch', 'Kafka', 'RabbitMQ'],
        'Soft Skills': ['Problem Solving', 'Teamwork', 'Project Management', 'Agile', 'Communication', 'Data Analysis']
    }
    
    categorized_skills = {}
    for skill in all_found_skills:
        for category, skills_list in skill_categories.items():
            if skill in skills_list:
                if category not in categorized_skills:
                    categorized_skills[category] = []
                categorized_skills[category].append(skill)
                break
    
    # Get role recommendations based on extracted skills
    role_suggestions = career_system.recommend_roles(all_found_skills, top_n=3)
    
    # Format suggestions
    top_role = role_suggestions[0] if role_suggestions else None
    suggested_roles = [
        {
            'role': suggestion['role'],
            'match_percentage': suggestion['match_percentage'],
            'confidence': 'High' if suggestion['match_percentage'] > 70 else 'Medium' if suggestion['match_percentage'] > 50 else 'Low'
        }
        for suggestion in role_suggestions
    ]
    
    # Save resume to database
    try:
        resume = Resume(
            filename=original_filename,
            file_path=file_path,
            file_type=file_type,
            file_size=file_size,
            text_content=content[:10000],  # Store first 10000 chars to avoid DB bloat
            extracted_skills=json.dumps(all_found_skills),
            categorized_skills=json.dumps(categorized_skills),
            suggested_roles=json.dumps(suggested_roles),
            top_match_role=top_role['role'] if top_role else None,
            top_match_percentage=top_role['match_percentage'] if top_role else None,
            processed=True,
            upload_date=datetime.utcnow()
        )
        
        db.session.add(resume)
        db.session.commit()
        
        resume_id = resume.id
        print(f"✓ Resume saved to database with ID: {resume_id}")
        
    except Exception as e:
        print(f"Error saving resume to database: {str(e)}")
        db.session.rollback()
        resume_id = None
    
    response_data = {
        'resume_id': resume_id,
        'skills': all_found_skills,
        'count': len(all_found_skills),
        'file_type': file_type,
        'categorized_skills': categorized_skills,
        'suggested_roles': suggested_roles,
        'top_match': {
            'role': top_role['role'] if top_role else None,
            'match_percentage': top_role['match_percentage'] if top_role else 0,
            'message': f"Based on your resume, you're best suited for {top_role['role']} role with {top_role['match_percentage']}% match!" if top_role else "Upload your resume to get role suggestions"
        },
        'saved_to_database': resume_id is not None
    }
    
    return jsonify(response_data)

@app.route('/api/roles', methods=['GET'])
def get_all_roles():
    """Get all available job roles"""
    roles = df_encoded['role'].unique().tolist()
    return jsonify({'roles': roles})

@app.route('/api/analyze', methods=['POST'])
def analyze_skills():
    """Analyze skill gap for a specific role"""
    data = request.json
    user_skills = data.get('skills', [])
    target_role = data.get('target_role')
    
    if not user_skills:
        return jsonify({'error': 'No skills provided'}), 400
    
    if not target_role:
        return jsonify({'error': 'No target role provided'}), 400
    
    # Analyze skill gap
    result = career_system.analyze_skill_gap(user_skills, target_role)
    
    if result is None:
        return jsonify({'error': 'Invalid target role'}), 400
    
    return jsonify(result)

@app.route('/api/recommend', methods=['POST'])
def recommend_roles():
    """Recommend suitable roles using Random Forest ML Model"""
    data = request.json
    user_skills = data.get('skills', [])
    top_n = data.get('top_n', 5)
    
    if not user_skills:
        return jsonify({'error': 'No skills provided'}), 400
    
    # Use ML model for predictions
    recommendations = career_system.recommend_roles(user_skills, top_n)
    
    return jsonify({
        'recommendations': recommendations,
        'model_used': 'Random Forest Classifier',
        'total_skills_provided': len(user_skills)
    })

@app.route('/api/career-guidance', methods=['POST'])
def career_guidance():
    """Generate AI-powered career guidance using LLaMA"""
    data = request.json
    user_skills = data.get('skills', [])
    target_role = data.get('target_role')
    missing_skills = data.get('missing_skills', [])
    
    if not user_skills or not target_role:
        return jsonify({'error': 'Missing required parameters'}), 400
    
    # Get skill gap analysis
    analysis = career_system.analyze_skill_gap(user_skills, target_role)
    
    if analysis is None:
        return jsonify({'error': 'Invalid target role'}), 400
    
    matching_skills = analysis.get('matching_skills', [])
    missing_skills = analysis.get('missing_skills', []) if not missing_skills else missing_skills
    
    # Check if LLaMA is available
    ollama_available = llama_ai.check_ollama_status()
    
    # Generate skill gap explanation
    skill_gap_explanation = llama_ai.explain_skill_gap(
        user_skills, target_role, missing_skills, matching_skills
    )
    
    # Generate learning roadmap
    learning_roadmap = llama_ai.generate_learning_roadmap(
        user_skills, target_role, missing_skills
    )
    
    response = {
        'target_role': target_role,
        'match_percentage': analysis['match_percentage'],
        'skills_analysis': {
            'total_skills_provided': len(user_skills),
            'matching_skills_count': len(matching_skills),
            'missing_skills_count': len(missing_skills),
            'matching_skills': matching_skills[:10],  # Top 10
            'missing_skills': missing_skills[:10]  # Top 10
        },
        'skill_gap_explanation': skill_gap_explanation,
        'learning_roadmap': learning_roadmap,
        'ai_powered': ollama_available,
        'llama_model': LLAMA_MODEL if ollama_available else 'Fallback mode'
    }
    
    return jsonify(response)

@app.route('/api/predict-career', methods=['POST'])
def predict_career():
    """Predict career role using Random Forest Classifier"""
    data = request.json
    user_skills = data.get('skills', [])
    
    if not user_skills:
        return jsonify({'error': 'No skills provided'}), 400
    
    # Get predictions from ML model
    predictions = career_system.predict_career(user_skills)
    
    return jsonify({
        'predictions': predictions,
        'model': 'Random Forest Classifier',
        'top_prediction': predictions[0] if predictions else None,
        'skills_analyzed': len(user_skills)
    })

# ========== Resume Management Endpoints ==========

@app.route('/api/resumes', methods=['GET'])
def get_all_resumes():
    """Get all uploaded resumes"""
    try:
        limit = request.args.get('limit', 50, type=int)
        resumes = Resume.get_all_resumes(limit=limit)
        
        return jsonify({
            'resumes': [resume.to_dict() for resume in resumes],
            'count': len(resumes),
            'total_in_db': Resume.query.count()
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/resumes/<int:resume_id>', methods=['GET'])
def get_resume(resume_id):
    """Get a specific resume by ID"""
    try:
        resume = Resume.get_resume_by_id(resume_id)
        
        if not resume:
            return jsonify({'error': 'Resume not found'}), 404
        
        return jsonify(resume.to_dict())
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/resumes/<int:resume_id>/download', methods=['GET'])
def download_resume(resume_id):
    """Download the original resume file"""
    try:
        resume = Resume.get_resume_by_id(resume_id)
        
        if not resume:
            return jsonify({'error': 'Resume not found'}), 404
        
        if not resume.file_path or not os.path.exists(resume.file_path):
            return jsonify({'error': 'Resume file not found on disk'}), 404
        
        return send_file(
            resume.file_path,
            as_attachment=True,
            download_name=resume.filename
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/resumes/<int:resume_id>', methods=['DELETE'])
def delete_resume(resume_id):
    """Delete a resume from database and disk"""
    try:
        resume = Resume.get_resume_by_id(resume_id)
        
        if not resume:
            return jsonify({'error': 'Resume not found'}), 404
        
        # Delete file from disk
        if resume.file_path and os.path.exists(resume.file_path):
            try:
                os.remove(resume.file_path)
                print(f"✓ Deleted file: {resume.file_path}")
            except Exception as e:
                print(f"Warning: Could not delete file: {str(e)}")
        
        # Delete from database
        Resume.delete_resume(resume_id)
        
        return jsonify({
            'message': 'Resume deleted successfully',
            'resume_id': resume_id
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/resumes/stats', methods=['GET'])
def get_resume_stats():
    """Get statistics about uploaded resumes"""
    try:
        total_resumes = Resume.query.count()
        total_size = db.session.query(db.func.sum(Resume.file_size)).scalar() or 0
        
        # Get most common skills across all resumes
        all_skills_data = Resume.query.with_entities(Resume.extracted_skills).all()
        skill_counter = {}
        
        for (skills_json,) in all_skills_data:
            if skills_json:
                skills = json.loads(skills_json)
                for skill in skills:
                    skill_counter[skill] = skill_counter.get(skill, 0) + 1
        
        # Get top 10 most common skills
        top_skills = sorted(skill_counter.items(), key=lambda x: x[1], reverse=True)[:10]
        
        # Get role distribution
        role_distribution = {}
        roles_data = Resume.query.with_entities(Resume.top_match_role).all()
        for (role,) in roles_data:
            if role:
                role_distribution[role] = role_distribution.get(role, 0) + 1
        
        return jsonify({
            'total_resumes': total_resumes,
            'total_size_bytes': total_size,
            'total_size_mb': round(total_size / (1024 * 1024), 2),
            'top_skills': [{'skill': skill, 'count': count} for skill, count in top_skills],
            'role_distribution': role_distribution,
            'avg_skills_per_resume': round(sum(skill_counter.values()) / total_resumes, 2) if total_resumes > 0 else 0
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/skill-gap-analyses', methods=['GET'])
def get_all_analyses():
    """Get all skill gap analyses"""
    try:
        limit = request.args.get('limit', 50, type=int)
        analyses = SkillGapAnalysis.query.order_by(
            SkillGapAnalysis.analysis_date.desc()
        ).limit(limit).all()
        
        return jsonify({
            'analyses': [analysis.to_dict() for analysis in analyses],
            'count': len(analyses)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/skill-gap-analyses', methods=['POST'])
def save_skill_gap_analysis():
    """Save a skill gap analysis to database"""
    try:
        data = request.json
        
        analysis = SkillGapAnalysis(
            resume_id=data.get('resume_id'),
            user_skills=json.dumps(data.get('user_skills', [])),
            target_role=data.get('target_role'),
            match_percentage=data.get('match_percentage', 0),
            missing_skills=json.dumps(data.get('missing_skills', [])),
            matching_skills=json.dumps(data.get('matching_skills', [])),
            recommended_skills=json.dumps(data.get('recommended_skills', [])),
            ai_explanation=data.get('ai_explanation'),
            learning_roadmap=json.dumps(data.get('learning_roadmap'))
        )
        
        db.session.add(analysis)
        db.session.commit()
        
        return jsonify({
            'message': 'Analysis saved successfully',
            'analysis_id': analysis.id,
            'analysis': analysis.to_dict()
        }), 201
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000, host='0.0.0.0')
